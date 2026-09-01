"""Extract remaining 4 gammes via Groq LLM. Run: python extract_remaining.py"""
import json, os, time
from pptx import Presentation
from groq import Groq
import os

client = Groq(api_key=os.environ.get('GROQ_API_KEY', ''))
RAW = 'data/raw'

EXTRACT = (
    "Extract product data from this VITAL SA pharmaceutical slide. "
    "Return ONLY valid JSON (no markdown, no thinking):\n"
    '{"name":"product name","presentation":"Sirop/Gelules/Creme/etc","packaging":"",'
    '"age_range":"Adulte/Enfant or null","indications":["what it treats"],'
    '"composition":[{"quantity":"20 mg","ingredient":"Ingredient Name"}],'
    '"posologie":{"note":"","table":null}}\n'
    "Ignore headers like INDICATIONS, COMPOSITION, POSOLOGIE.\n\n"
    "RAW TEXT:\n---\n{raw}\n---\n\nJSON:"
)

with open('data/processed/products_llm.json', 'r', encoding='utf-8') as f:
    existing = json.load(f)

new_products = []
for fname in ['Gamme TC2000 et SPS ARNICA.pptx', 'Gamme Tidol.pptx', 'Gamme UNIDERM.pptx', 'Gamme VITONIC.pptx']:
    gamme = fname.replace('Gamme ', '').replace('.pptx', '')
    prs = Presentation(os.path.join(RAW, fname))
    count = 0
    for slide in prs.slides:
        texts = [s.text.strip() for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
        raw = chr(10).join(texts)
        if not raw or len(raw) < 20:
            continue
        prompt = EXTRACT.replace('{raw}', raw[:2500])
        for attempt in range(3):
            try:
                resp = client.chat.completions.create(
                    model='qwen/qwen3.8-27b',
                    messages=[
                        {'role': 'system', 'content': 'Return ONLY valid JSON.'},
                        {'role': 'user', 'content': prompt}
                    ],
                    temperature=0.1, max_tokens=1200
                )
                content = (resp.choices[0].message.content or '').strip()
                if chr(8203) in content:
                    content = content[content.rfind('{'):]
                if content.startswith('```'):
                    content = content.split(chr(10), 1)[1]
                if content.endswith('```'):
                    content = content[:-3]
                data = json.loads(content.strip())
                if not data.get('name'):
                    data['name'] = gamme
                data['gamme'] = gamme
                new_products.append(data)
                count += 1
                break
            except json.JSONDecodeError:
                pass
            except Exception as e:
                if '429' in str(e):
                    time.sleep(15)
                    continue
                break
        time.sleep(6)
    print(f'{gamme}: {count} products')

all_products = existing + new_products
with open('data/processed/products_llm.json', 'w', encoding='utf-8') as f:
    json.dump(all_products, f, ensure_ascii=False, indent=2)

gammes = {}
for p in all_products:
    g = p.get('gamme', '?')
    gammes[g] = gammes.get(g, 0) + 1
print(f'\nTotal: {len(all_products)} products across {len(gammes)} gammes')
for g, c in sorted(gammes.items()):
    print(f'  {g}: {c}')
