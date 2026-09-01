"""
ALIA Avatar - Deep PPTX Parser v2
Extracts structured product data from VITAL SA gamme PowerPoint files.
Handles the pattern: Slide 1 = gamme title, Slide 2+ = individual product slides
with product name embedded in the same slide as indications/composition/posologie.
"""
import os
import re
import json
from typing import List, Dict, Any, Optional, Tuple
from loguru import logger

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PPTXProductParser:
    """Parse VITAL SA gamme PPTX files into structured JSON."""

    def __init__(self, catalog_dir: str):
        self.catalog_dir = catalog_dir

    def parse_all(self) -> List[Dict[str, Any]]:
        """Parse all gamme PPTX files."""
        results = []
        if not os.path.exists(self.catalog_dir):
            logger.warning(f"Catalog directory not found: {self.catalog_dir}")
            return results

        pptx_files = sorted([f for f in os.listdir(self.catalog_dir) if f.endswith(".pptx")])
        logger.info(f"Found {len(pptx_files)} PPTX files to parse")

        for filename in pptx_files:
            gamme_name = filename.replace("Gamme ", "").replace(".pptx", "")
            filepath = os.path.join(self.catalog_dir, filename)
            gamme_data = self.parse_gamme(gamme_name, filepath)
            if gamme_data and gamme_data["products"]:
                results.append(gamme_data)
                logger.info(f"  {gamme_name}: {len(gamme_data['products'])} products extracted")

        return results

    def parse_gamme(self, gamme_name: str, filepath: str) -> Optional[Dict[str, Any]]:
        """Parse a single gamme PPTX."""
        if not HAS_PPTX:
            return None

        try:
            prs = Presentation(filepath)
        except Exception as e:
            logger.error(f"Failed to open {filepath}: {e}")
            return None

        gamme_data = {
            "gamme": gamme_name,
            "source_file": os.path.basename(filepath),
            "products": [],
        }

        # Collect all slides with their text, tables, and shapes
        slides_data = []
        for slide in prs.slides:
            text = self._extract_text(slide)
            tables = self._extract_tables(slide)
            product_name_shape = self._extract_product_name_from_shapes(slide)
            slides_data.append({
                "text": text,
                "tables": tables,
                "product_name_shape": product_name_shape,
                "raw_shapes": self._get_shape_names(slide),
                "slide": slide,  # Keep reference for shape access
            })

        # Detect product boundaries
        products = self._split_into_products(gamme_name, slides_data)
        gamme_data["products"] = products

        return gamme_data

    def _get_shape_names(self, slide) -> List[str]:
        """Get shape names/types for debugging."""
        names = []
        for shape in slide.shapes:
            names.append(f"{shape.shape_type}:{getattr(shape, 'name', '?')}")
        return names

    def _extract_text(self, slide) -> str:
        """Extract all text from a slide."""
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        return "\n".join(texts)

    def _extract_tables(self, slide) -> List[List[List[str]]]:
        """Extract tables from slide."""
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                table_data = []
                for row in list(shape.table.rows):
                    cells = [cell.text.strip() for cell in list(row.cells)]
                    table_data.append(cells)
                tables.append(table_data)
        return tables

    def _split_into_products(self, gamme_name: str, slides_data: List[Dict]) -> List[Dict[str, Any]]:
        """Split slides into individual products.
        
        Each slide typically = one product (especially in PEDIKIDS-style gammes).
        The product name is found in specific shapes or as the first short text line.
        """
        products = []

        # Skip first slide if it's just the gamme title
        start_idx = 0
        if slides_data:
            first_text = slides_data[0]["text"].strip()
            if first_text.upper() == gamme_name.upper() or len(first_text.split("\n")) <= 2:
                start_idx = 1

        product_count = 0

        for i in range(start_idx, len(slides_data)):
            slide = slides_data[i]
            text = slide["text"]
            tables = slide["tables"]
            shape_name = slide.get("product_name_shape")

            # Skip empty slides
            if not text.strip() and not tables:
                continue

            # Build product from this single slide
            product = self._build_product_from_slide(gamme_name, text, tables, product_count, shape_name)
            if product:
                products.append(product)
                product_count += 1

        return products

    def _build_product_from_slide(self, gamme_name: str, text: str, tables: List, product_idx: int, shape_name: Optional[str] = None) -> Optional[Dict[str, Any]]:
        """Build a product from a single slide's data."""
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        if not lines:
            return None

        # Extract product name - try multiple strategies
        name = ""
        # Strategy 1: From shape name (e.g. "APITOU N°1" in PEDIKIDS)
        if shape_name:
            name = shape_name
        # Strategy 2: From text lines
        if not name:
            name = self._extract_product_name(lines, gamme_name)
        if not name:
            name = f"{gamme_name} Product {product_idx + 1}"

        return {
            "name": name,
            "gamme": gamme_name,
            "presentation": self._extract_presentation(lines),
            "packaging": self._extract_packaging(lines),
            "age_range": self._extract_age_range(lines),
            "indications": self._extract_indications(text),
            "composition": self._extract_composition(text),
            "posologie": self._extract_posologie(tables),
            "raw_text": text.strip(),
        }

    def _detect_new_product(self, lines: List[str], is_first: bool) -> bool:
        """Detect if a slide starts a new product."""
        if is_first:
            return True

        if not lines:
            return False

        # Content keywords indicate this is NOT a new product slide
        content_keywords = [
            "INDICATION", "COMPOSITION", "POSOLOGIE", "CONTRAINDICATION",
            "PRÉCAUTION", "PRECAUTION", "CONSEIL", "AVIS", "MODE D'EMPLOI",
            "CONTRE", "MATIN", "MIDI", "SOIR", "AGE", "ÉGE",
            "BÉNÉFICE", "AVANTAGE", "EFFET", "ACTION", "TRAITMENT",
        ]
        joined_upper = " ".join(lines).upper()
        if any(kw in joined_upper for kw in content_keywords):
            return False

        # A product name slide has 1-3 short lines
        if len(lines) <= 3 and all(len(l) < 50 for l in lines):
            return True

        return False

    def _build_product(self, gamme_name: str, text: str, tables: List, product_idx: int) -> Optional[Dict[str, Any]]:
        """Build a structured product dict from collected text and tables."""
        lines = [l.strip() for l in text.split("\n") if l.strip()]

        if not lines:
            return None

        # Extract product name (first meaningful short line)
        name = self._extract_product_name(lines, gamme_name)
        if not name:
            name = f"{gamme_name} Product {product_idx + 1}"

        product = {
            "name": name,
            "gamme": gamme_name,
            "presentation": self._extract_presentation(lines),
            "packaging": self._extract_packaging(lines),
            "age_range": self._extract_age_range(lines),
            "indications": self._extract_indications(text),
            "composition": self._extract_composition(text),
            "posologie": self._extract_posologie(tables),
            "raw_text": text.strip(),
        }

        return product

    def _extract_product_name_from_shapes(self, slide) -> Optional[str]:
        """Extract product name from specific shape names (e.g. 'ZoneTexte 20')."""
        for shape in slide.shapes:
            name = getattr(shape, "name", "")
            # PEDIKIDS uses 'ZoneTexte 20' for product names
            if "ZoneTexte 20" in name or "ZoneTexte" in name:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Validate it looks like a product name (short, no content keywords)
                    if len(text) < 50 and not any(kw in text.upper() for kw in [
                        "INDICATION", "COMPOSITION", "POSOLOGIE", "ML", "MG"
                    ]):
                        return text
        return None

    def _extract_product_name(self, lines: List[str], gamme_name: str) -> str:
        """Extract the product name from text lines."""
        skip_words = [
            "INDICATION", "COMPOSITION", "POSOLOGIE", "CONTRAINDICATION",
            "PRÉCAUTION", "PRECAUTION", "CONSEIL", "AVIS", "AGE",
            "MATIN", "MIDI", "SOIR", "BÉNÉFICE", "AVANTAGE",
        ]

        for line in lines[:5]:
            upper = line.upper()
            if any(sw in upper for sw in skip_words):
                continue
            if len(line) < 3 or len(line) > 50:
                continue
            if re.match(r"^[\d\s,\.]+(mg|µg|g|ml|%)?$", line, re.IGNORECASE):
                continue
            if line.upper() == gamme_name.upper():
                continue
            if re.match(r"^\d+\s*(ML|Gélules?|Sachets?|Comprimés?|g)$", line, re.IGNORECASE):
                continue
            return line.strip()

        return ""

    def _extract_indications(self, text: str) -> List[str]:
        """Extract indications from full slide text."""
        indications = []
        # Find INDICATIONS section
        match = re.search(r"INDICATIONS?\s*[:\n](.*?)(?:COMPOSITION|POSOLOGIE|CONTRAINDICATION|$)",
                         text, re.IGNORECASE | re.DOTALL)
        if match:
            section = match.group(1)
            for line in section.split("\n"):
                clean = line.strip().lstrip("- •●◦▸▹►*")
                clean = clean.strip()
                if clean and len(clean) > 3 and len(clean) < 300:
                    # Skip composition/posologie leaks
                    if not re.match(r"^[\d,\.]+\s*(mg|µg|g|ml)", clean, re.IGNORECASE):
                        indications.append(clean)
        return indications

    def _extract_composition(self, text: str) -> List[Dict[str, str]]:
        """Extract composition ingredients."""
        composition = []
        match = re.search(r"COMPOSITION\s*(?:PAR\s*\d+\s*(?:ml|g|sachet|gélule))?\s*[:\n](.*?)(?:POSOLOGIE|INDICATION|CONTRAINDICATION|$)",
                         text, re.IGNORECASE | re.DOTALL)
        if match:
            section = match.group(1)
            for line in section.split("\n"):
                line = line.strip()
                if not line or len(line) < 3:
                    continue
                # Skip product names / gamme headers
                skip_patterns = [
                    r"^(FERBIOTIC|CALMOSS|PEDIAKIDS|COSMOPHARMA|BACTOL|HEALTHCARE|HYDRA|"
                    r"LABORATOIRE VITAL|MINCILIGNE|MINCIVIT|OLIGOVIT|OMEVIE|PHYTOL|"
                    r"PHYTOPHANE|PHYTOTHERA|PLANTHERAPIE|TC2000|TIDOL|UNIDERM|VITONIC|VITOSINE)",
                    r"^\d+\s*(ML|Gélules?|Sachets?|Comprimés?)$",
                    r"^(COMPOSITION|POSOLOGIE|INDICATION)",
                ]
                if any(re.match(p, line, re.IGNORECASE) for p in skip_patterns):
                    continue

                # Parse ingredient: "80 mg Vitamine C" or "25 mg de Histidine"
                parsed = self._parse_ingredient_line(line)
                if parsed:
                    composition.append(parsed)

        return composition

    def _parse_ingredient_line(self, line: str) -> Optional[Dict[str, str]]:
        """Parse a single ingredient line."""
        # Pattern: "quantity ingredient" e.g. "80 mg Vitamine C"
        match = re.match(r"([\d,\.]+\s*(?:mg|µg|g|ml|UI|%))\s+(?:de\s+)?(.+)", line, re.IGNORECASE)
        if match:
            return {"quantity": match.group(1).strip(), "ingredient": match.group(2).strip()}

        # Pattern: just ingredient name
        if len(line) > 3 and not re.match(r"^\d", line):
            return {"quantity": "", "ingredient": line}

        return None

    def _extract_posologie(self, tables: List[List[List[str]]]) -> Dict[str, Any]:
        """Extract posologie from tables."""
        for table in tables:
            if not table:
                continue
            header = " ".join(table[0]).upper()
            if any(w in header for w in ["AGE", "ÂGE", "MATIN", "MIDI", "SOIR", "DOSE", "ÉGE"]):
                return {
                    "headers": table[0],
                    "rows": table[1:],
                }
        return {"headers": [], "rows": []}

    def _extract_packaging(self, lines: List[str]) -> str:
        """Extract packaging info."""
        for line in lines:
            match = re.search(r"(\d+\s*(?:Gélules?|Sachets?|Comprimés?|ML|mg|g|Flacon|Tube|Boîte|Spray))", line, re.IGNORECASE)
            if match:
                return match.group(1).strip()
        return ""

    def _extract_age_range(self, lines: List[str]) -> str:
        """Extract age range."""
        for line in lines:
            upper = line.upper()
            if "DÈS LA NAISSANCE" in upper:
                return "Dès la naissance"
            if any(w in upper for w in ["ADULTE", "ADULTES"]):
                return "Adulte"
            if any(w in upper for w in ["ENFANT", "PÉDIATRIQUE", "BÉBÉ"]):
                match = re.search(r"(\d+\s*(?:à|ans?|mois))", line, re.IGNORECASE)
                if match:
                    return f"Enfant {match.group(0)}"
                return "Enfant"
        return ""

    def _extract_presentation(self, lines: List[str]) -> str:
        """Extract presentation type."""
        pres_keywords = ["Sirop", "Gélules", "Gélule", "Sachets", "Sachet",
                        "Comprimé à sucer", "Comprimés", "Comprimé", "Spray",
                        "Crème", "Gel", "Poudre", "Solution", "Spray Gorge"]
        for line in lines:
            for kw in pres_keywords:
                if kw.lower() in line.lower():
                    return kw
        return ""

    def save_json(self, data: Any, output_path: str):
        """Save as JSON."""
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        logger.info(f"Saved: {output_path}")

    def to_flat_products(self, gamme_data: List[Dict]) -> List[Dict[str, Any]]:
        """Flatten into individual product records."""
        flat = []
        for gamme in gamme_data:
            for product in gamme["products"]:
                flat.append({
                    "name": product["name"],
                    "gamme": gamme["gamme"],
                    "presentation": product.get("presentation", ""),
                    "packaging": product.get("packaging", ""),
                    "age_range": product.get("age_range", ""),
                    "indications": product.get("indications", []),
                    "composition": product.get("composition", []),
                    "posologie": product.get("posologie", {}),
                    "raw_text": product.get("raw_text", ""),
                })
        return flat


if __name__ == "__main__":
    import sys
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))

    CATALOG = os.path.join(os.path.dirname(__file__), "..", "..", "..", "..", "Students", "Data vital", "catalogue", "Catalogue vital")
    OUTPUT = os.path.join(os.path.dirname(__file__), "..", "..", "data", "processed")

    parser = PPTXProductParser(CATALOG)
    gamme_data = parser.parse_all()

    # Save gamme-structured data
    parser.save_json(gamme_data, os.path.join(OUTPUT, "gamme_parsed.json"))

    # Save flat products
    flat = parser.to_flat_products(gamme_data)
    parser.save_json(flat, os.path.join(OUTPUT, "products_parsed.json"))

    total = sum(len(g["products"]) for g in gamme_data)
    print(f"Parsed {len(gamme_data)} gammes -> {total} products")
    for g in gamme_data:
        for p in g["products"]:
            ind_count = len(p["indications"])
            comp_count = len(p["composition"])
            poso = "Yes" if p["posologie"]["headers"] else "No"
            print(f"  {p['name']}: {ind_count} indications, {comp_count} ingredients, poso={poso}, pres={p['presentation']}")
