"""
ALIA Avatar - Product Catalog Parser
Parses VITAL SA product catalog PPTX files into structured data.
"""
import os
import json
import re
from typing import List, Dict, Any, Optional
from loguru import logger

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed. PPTX parsing unavailable.")

from app.ai.rag import rag_pipeline


# ──────────────────────────────────────────────
# Product Catalog Definitions
# ──────────────────────────────────────────────

PRODUCT_CATALOG = {
    "BACTOL": {
        "name": "BACTOL",
        "category": "Antiseptique / Hygiène",
        "therapeutic_area": "Hygiène buccale",
        "specialties": ["MG", "ORL", "Dentiste"],
    },
    "CALMOSS": {
        "name": "CALMOSS",
        "category": "Phytoterapie",
        "therapeutic_area": "Calme / Sommeil",
        "specialties": ["MG", "Neurologue"],
    },
    "COSMOPHARMA": {
        "name": "COSMOPHARMA",
        "category": "Cosmetique / Dermatologie",
        "therapeutic_area": "Soins dermatologiques",
        "specialties": ["Dermatologue", "MG"],
    },
    "FERBIOTIC": {
        "name": "FERBIOTIC",
        "category": "Fer / Energie",
        "therapeutic_area": "Carence en fer",
        "specialties": ["MG", "Gynecologue"],
    },
    "HEALTHCARE": {
        "name": "HEALTHCARE",
        "category": "Sante generale",
        "therapeutic_area": "Bien-etre",
        "specialties": ["MG"],
    },
    "HYDRA": {
        "name": "HYDRA",
        "category": "Hydratation",
        "therapeutic_area": "Hydratation cutanee",
        "specialties": ["Dermatologue", "MG"],
    },
    "LABORATOIRE VITAL": {
        "name": "LABORATOIRE VITAL",
        "category": "Pharmaceutique general",
        "therapeutic_area": "Polyvalent",
        "specialties": ["MG", "Specialiste"],
    },
    "MINCILIGNE": {
        "name": "MINCILIGNE",
        "category": "Perte de poids",
        "therapeutic_area": "Poids / Metabolisme",
        "specialties": ["MG", "Nutritionniste"],
    },
    "MINCIVIT": {
        "name": "MINCIVIT",
        "category": "Complement alimentaire",
        "therapeutic_area": "Poids / Energie",
        "specialties": ["MG"],
    },
    "OLIGOVIT": {
        "name": "OLIGOVIT",
        "category": "Vitamines / Mineraux",
        "therapeutic_area": "Carences vitaminiques",
        "specialties": ["MG"],
    },
    "OMEVIE": {
        "name": "OMEVIE",
        "category": "Omega-3",
        "therapeutic_area": "Cardiovasculaire / Bien-etre",
        "specialties": ["MG", "Cardiologue"],
    },
    "PEDIAKIDS": {
        "name": "PEDIAKIDS",
        "category": "Pediatrie",
        "therapeutic_area": "Sante infantile",
        "specialties": ["Pediatre", "MG"],
    },
    "PHYTOL": {
        "name": "PHYTOL",
        "category": "Phytoterapie",
        "therapeutic_area": "Bien-etre / Naturel",
        "specialties": ["MG"],
    },
    "PHYTOPHANE": {
        "name": "PHYTOPHANE",
        "category": "Phytoterapie",
        "therapeutic_area": "Cheveux / Ongles",
        "specialties": ["Dermatologue", "MG"],
    },
    "PHYTOTHERA": {
        "name": "PHYTOTHERA",
        "category": "Phytoterapie / Medical",
        "therapeutic_area": "Therapeutique naturelle",
        "specialties": ["MG", "Specialiste"],
    },
    "PLANTHERAPIE": {
        "name": "PLANTHERAPIE",
        "category": "Phytoterapie",
        "therapeutic_area": "Therapeutique par les plantes",
        "specialties": ["MG"],
    },
    "TC2000": {
        "name": "TC2000 et SPS ARNICA",
        "category": "Homeopathie / Phyto",
        "therapeutic_area": "Douleurs / Contusions",
        "specialties": ["MG", "Rhumatologue"],
    },
    "Tidol": {
        "name": "Tidol",
        "category": "Anti-douleur",
        "therapeutic_area": "Douleurs legeres",
        "specialties": ["MG"],
    },
    "UNIDERM": {
        "name": "UNIDERM",
        "category": "Dermatologie",
        "therapeutic_area": "Soins dermatologiques",
        "specialties": ["Dermatologue", "MG"],
    },
    "VITONIC": {
        "name": "VITONIC",
        "category": "Vitamines",
        "therapeutic_area": "Energie / Metabolisme",
        "specialties": ["MG", "Gynecologue"],
    },
    "Vitosine": {
        "name": "Vitosine",
        "category": "Complement alimentaire",
        "therapeutic_area": "Bien-etre general",
        "specialties": ["MG"],
    },
}


class ProductCatalogParser:
    """Parse product catalog PPTX files into structured data."""

    def __init__(self, catalog_dir: str):
        self.catalog_dir = catalog_dir

    def parse_all_gammes(self) -> List[Dict[str, Any]]:
        """Parse all gamme PPTX files in the catalog directory."""
        products = []

        if not os.path.exists(self.catalog_dir):
            logger.warning(f"Catalog directory not found: {self.catalog_dir}")
            return self._get_default_products()

        for filename in os.listdir(self.catalog_dir):
            if filename.endswith(".pptx") and filename.startswith("Gamme"):
                gamme_name = filename.replace("Gamme ", "").replace(".pptx", "")
                product_data = self.parse_gamme(gamme_name, os.path.join(self.catalog_dir, filename))
                if product_data:
                    products.append(product_data)

        if not products:
            logger.warning("No PPTX files found, using defaults")
            return self._get_default_products()

        return products

    def parse_gamme(self, gamme_name: str, filepath: str) -> Optional[Dict[str, Any]]:
        """Parse a single gamme PPTX file."""
        if not HAS_PPTX:
            return self._get_default_product(gamme_name)

        try:
            prs = Presentation(filepath)
            slides_text = []
            products_in_gamme = []

            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides_text.append("\n".join(slide_text))

            full_text = "\n\n".join(slides_text)

            # Extract product information
            catalog_info = PRODUCT_CATALOG.get(gamme_name, {})

            # Parse products from slides
            product_data = {
                "name": gamme_name,
                "category": catalog_info.get("category", "Non categorise"),
                "therapeutic_area": catalog_info.get("therapeutic_area", ""),
                "specialties": catalog_info.get("specialties", []),
                "gamme": gamme_name,
                "slides_content": full_text[:2000],  # First 2000 chars for context
                "key_benefits": self._extract_benefits(full_text),
                "dosage": self._extract_dosage(full_text),
                "precautions": self._extract_precautions(full_text),
                "composition": self._extract_composition(full_text),
                "evidence": [],
                "target_profiles": self._extract_target_profiles(full_text),
            }

            logger.info(f"Parsed gamme: {gamme_name} ({len(slides_text)} slides)")
            return product_data

        except Exception as e:
            logger.error(f"Failed to parse {filepath}: {e}")
            return self._get_default_product(gamme_name)

    def _extract_benefits(self, text: str) -> List[str]:
        """Extract key benefits from text."""
        benefits = []
        benefit_patterns = [
            r"b[eé]n[eé]fice[s]?\s*:?\s*(.+?)(?:\n|$)",
            r"avantage[s]?\s*:?\s*(.+?)(?:\n|$)",
            r"efficacit[eé]\s*:?\s*(.+?)(?:\n|$)",
        ]
        for pattern in benefit_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            benefits.extend([m.strip() for m in matches if m.strip()])
        return benefits[:5]

    def _extract_dosage(self, text: str) -> str:
        """Extract dosage information."""
        dosage_patterns = [
            r"dosage\s*:?\s*(.+?)(?:\n|$)",
            r"posologie\s*:?\s*(.+?)(?:\n|$)",
            r"(\d+\s*(?:mg|g|ml|cp|gelule|sachet).+?)(?:\n|$)",
        ]
        for pattern in dosage_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1).strip()[:200]
        return ""

    def _extract_precautions(self, text: str) -> List[str]:
        """Extract precautions."""
        precautions = []
        prec_patterns = [
            r"pr[eé]caution[s]?\s*:?\s*(.+?)(?:\n|$)",
            r"contre-indication[s]?\s*:?\s*(.+?)(?:\n|$)",
            r"avertissement[s]?\s*:?\s*(.+?)(?:\n|$)",
        ]
        for pattern in prec_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            precautions.extend([m.strip() for m in matches if m.strip()])
        return precautions[:5]

    def _extract_composition(self, text: str) -> str:
        """Extract composition."""
        comp_patterns = [
            r"composition\s*:?\s*(.+?)(?:\n|$)",
            r"ingr[eé]dient[s]?\s*:?\s*(.+?)(?:\n|$)",
        ]
        for pattern in comp_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1).strip()[:200]
        return ""

    def _extract_target_profiles(self, text: str) -> List[str]:
        """Extract target patient profiles."""
        profiles = []
        if any(w in text.lower() for w in ["pediatrie", "enfant", "bebe"]):
            profiles.append("Enfants")
        if any(w in text.lower() for w in ["adulte", "senior"]):
            profiles.append("Adultes")
        if any(w in text.lower() for w in ["enceinte", "allaitement", "grossesse"]):
            profiles.append("Femmes enceintes/allaitantes")
        if any(w in text.lower() for w in ["sportif", "athlete"]):
            profiles.append("Sportifs")
        return profiles

    def _get_default_product(self, gamme_name: str) -> Dict[str, Any]:
        """Get default product data when parsing fails."""
        catalog_info = PRODUCT_CATALOG.get(gamme_name, {})
        return {
            "name": gamme_name,
            "category": catalog_info.get("category", "Non categorise"),
            "therapeutic_area": catalog_info.get("therapeutic_area", ""),
            "specialties": catalog_info.get("specialties", []),
            "gamme": gamme_name,
            "key_benefits": [],
            "dosage": "",
            "precautions": [],
            "composition": "",
            "evidence": [],
            "target_profiles": [],
        }

    def _get_default_products(self) -> List[Dict[str, Any]]:
        """Get default product list."""
        return [self._get_default_product(name) for name in PRODUCT_CATALOG.keys()]

    def ingest_all_products(self):
        """Parse and ingest all products into the RAG pipeline."""
        products = self.parse_all_gammes()

        for product in products:
            rag_pipeline.ingest_product(product)

        # Ingest visit process
        rag_pipeline.ingest_visit_process()

        logger.info(f"Ingested {len(products)} products into RAG pipeline")
        return products


# Global instance
product_parser = ProductCatalogParser(
    catalog_dir=os.path.join(os.path.dirname(__file__), "..", "..", "data", "raw")
)
