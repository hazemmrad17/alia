"""
Extract ALL 280 products from PPTX files using text parsing.
Run: python extract_all_products.py
"""
import os, sys, io, json, re
from pptx import Presentation
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

raw = 'data/raw'
all_products = []

SKIP_WORDS = {'INDICATIONS', 'COMPOSITION', 'POSOLOGIE', "CONSEILS D'UTILISATION"}

for fname in sorted(os.listdir(raw)):
    if not fname.endswith('.pptx'):
        continue
    gamme = fname.replace('Gamme ', '').replace('.pptx', '')
    fpath = os.path.join(raw, fname)
    prs = Presentation(fpath)
    count = 0

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
        full_text = '\n'.join(texts)
        if len(full_text) < 40:
            continue

        lines = [l.strip() for l in full_text.split('\n') if l.strip()]

        # Find product name
        product_name = None
        for line in lines:
            upper = line.upper()
            if upper in SKIP_WORDS:
                continue
            if upper.startswith('COMPOSITION'):
                continue
            if line.isupper() and 3 < len(line) < 80:
                product_name = line
                break
        if not product_name:
            product_name = f'{gamme} Product {count + 1}'

        # Extract indications
        indications = []
        in_ind = False
        for line in lines:
            u = line.upper()
            if 'INDICATIONS' in u and len(line) < 30:
                in_ind = True
                continue
            if 'COMPOSITION' in u or 'POSOLOGIE' in u or 'CONSEILS' in u:
                in_ind = False
                continue
            if in_ind:
                clean = re.sub(r'^[-\u2022*]\s*', '', line).strip()
                if clean and len(clean) > 3:
                    indications.append(clean)

        # Extract composition
        composition = []
        in_comp = False
        for line in lines:
            u = line.upper()
            if 'COMPOSITION' in u:
                in_comp = True
                continue
            if 'POSOLOGIE' in u or 'CONSEILS' in u:
                in_comp = False
                continue
            if in_comp:
                match = re.match(
                    r'([\d.,]+\s*(?:mg|\u00b5g|g|ml|UI)?)\s*(?:de\s+)?(.+)',
                    line, re.IGNORECASE
                )
                if match:
                    composition.append({
                        'quantity': match.group(1).strip(),
                        'ingredient': match.group(2).strip()
                    })
                elif len(line) > 3 and not line[0].isdigit():
                    composition.append({'quantity': '', 'ingredient': line})

        # Packaging
        packaging = ''
        for line in lines:
            m = re.search(
                r'\b(\d+\s*(?:g|gr|ml|MG|L|comprim\u00e9s?|g\u00e9lules?|sachets?|capsules?))\b',
                line, re.IGNORECASE
            )
            if m:
                packaging = m.group(1)
                break

        product = {
            'name': product_name,
            'gamme': gamme,
            'presentation': '',
            'packaging': packaging,
            'age_range': '',
            'indications': indications[:10],
            'composition': composition[:15],
            'posologie': {'note': '', 'table': None},
            'raw_text': full_text[:2000]
        }
        all_products.append(product)
        count += 1

    print(f'{gamme}: {count} products')

print(f'\nTotal: {len(all_products)} products across {len(set(p["gamme"] for p in all_products))} gammes')

with open('data/processed/products_full.json', 'w', encoding='utf-8') as f:
    json.dump(all_products, f, ensure_ascii=False, indent=2)
print('Saved to products_full.json')
